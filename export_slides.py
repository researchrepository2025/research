"""
Export Strategy Slides to PowerPoint and Word Documents
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import os

OUTPUT_DIR = "/Users/nicholaspate/Documents/01_Active/Corp_Strat/ezAI/exports"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Color scheme
COLORS = {
    'strategy_blue': RGBColor(0x1e, 0x40, 0xaf),
    'goals_green': RGBColor(0x05, 0x96, 0x69),
    'tactics_amber': RGBColor(0xd9, 0x77, 0x06),
    'correct_green': RGBColor(0x16, 0xa3, 0x4a),
    'incorrect_red': RGBColor(0xdc, 0x26, 0x26),
    'dark_gray': RGBColor(0x37, 0x41, 0x51),
    'light_gray': RGBColor(0xf3, 0xf4, 0xf6),
    'white': RGBColor(0xff, 0xff, 0xff),
}


def create_powerpoint():
    """Create PowerPoint with both strategy slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Strategy vs Goals vs Tactics
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategy vs. Goals vs. Tactics"
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.4))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Understanding the hierarchy of corporate planning"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    # Hierarchy bar
    hierarchy_box = slide1.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(6.333), Inches(0.4))
    tf = hierarchy_box.text_frame
    p = tf.paragraphs[0]
    p.text = "STRATEGY  -->  GOALS  -->  TACTICS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Three columns
    col_width = Inches(4.0)
    col_height = Inches(4.8)
    start_y = Inches(2.0)

    sections = [
        {
            'title': 'STRATEGY',
            'subtitle': 'Where & how to compete',
            'color': COLORS['strategy_blue'],
            'definition': 'An integrated set of choices that uniquely position an organization to create sustainable competitive advantage--including explicit trade-offs about what NOT to do.',
            'components': ['Where to Play', 'How to Win', 'Trade-offs'],
            'correct': '"Dominate mid-market manufacturing software through vertical specialization. We will NOT pursue enterprise accounts."',
            'incorrect': '"Be the best software company and grow revenue 20%."',
            'incorrect_note': 'No choices, no trade-offs'
        },
        {
            'title': 'GOALS',
            'subtitle': 'What outcomes to achieve',
            'color': COLORS['goals_green'],
            'definition': 'Specific, measurable outcomes with defined timeframes that translate strategic intent into concrete targets.',
            'components': ['Specific', 'Measurable', 'Achievable'],
            'correct': '"Acquire 500 new manufacturing customers in North America by December 2026."',
            'incorrect': '"Grow market share significantly."',
            'incorrect_note': 'Vague, no metrics, no deadline'
        },
        {
            'title': 'TACTICS',
            'subtitle': 'How to execute',
            'color': COLORS['tactics_amber'],
            'definition': 'Specific actions, initiatives, or methods employed to execute strategy and achieve goals--operational, short-term, and adjustable.',
            'components': ['Short-term', 'Actionable', 'Adjustable'],
            'correct': '"Sponsor NAM conference; hire 3 vertical sales reps; launch LinkedIn campaign."',
            'incorrect': '"Do more marketing and sales."',
            'incorrect_note': 'Not actionable, no specific activity'
        }
    ]

    for i, section in enumerate(sections):
        x = Inches(0.4 + i * 4.3)

        # Section header
        header_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, Inches(1.2), Inches(0.35))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = section['color']
        header_box.line.fill.background()
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = section['title']
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide1.shapes.add_textbox(x, Inches(2.4), Inches(4.0), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = section['subtitle']
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = COLORS['dark_gray']

        # Definition box
        def_box = slide1.shapes.add_textbox(x, Inches(2.7), Inches(4.0), Inches(1.0))
        tf = def_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "DEFINITION"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = section['color']
        p2 = tf.add_paragraph()
        p2.text = section['definition']
        p2.font.size = Pt(9)

        # Components
        comp_box = slide1.shapes.add_textbox(x, Inches(3.8), Inches(4.0), Inches(0.6))
        tf = comp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "KEY COMPONENTS"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = section['color']
        for comp in section['components']:
            p2 = tf.add_paragraph()
            p2.text = f"  {comp}"
            p2.font.size = Pt(9)

        # Correct example
        correct_box = slide1.shapes.add_textbox(x, Inches(4.7), Inches(4.0), Inches(0.8))
        tf = correct_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "CORRECT"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS['correct_green']
        p2 = tf.add_paragraph()
        p2.text = section['correct']
        p2.font.size = Pt(8)
        p2.font.italic = True

        # Incorrect example
        incorrect_box = slide1.shapes.add_textbox(x, Inches(5.6), Inches(4.0), Inches(0.8))
        tf = incorrect_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "INCORRECT"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = COLORS['incorrect_red']
        p2 = tf.add_paragraph()
        p2.text = section['incorrect']
        p2.font.size = Pt(8)
        p2.font.italic = True
        p3 = tf.add_paragraph()
        p3.text = f"  {section['incorrect_note']}"
        p3.font.size = Pt(7)
        p3.font.color.rgb = COLORS['incorrect_red']

    # Key insight footer
    insight_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6), Inches(13.333), Inches(0.4))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = COLORS['dark_gray']
    insight_box.line.fill.background()
    tf = insight_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategy defines what NOT to do  |  Goals are measurable targets  |  Tactics are swappable actions"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Sources
    source_box = slide1.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12), Inches(0.3))
    tf = source_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'Porter, M. "What Is Strategy?" HBR, 1996 | Martin, R. "The Big Lie of Strategic Planning" HBR, 2014 | McKinsey & Company'
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS['dark_gray']

    # =============================================
    # Slide 2: Strategy Creates Alignment
    # =============================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategy Creates Alignment--Its Absence Creates Chaos"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Left panel: WITHOUT STRATEGY
    left_panel = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.1), Inches(6.2), Inches(5.2))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0xfe, 0xe2, 0xe2)  # Light red
    left_panel.line.color.rgb = COLORS['incorrect_red']

    # Left header
    left_header = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = "WITHOUT STRATEGY"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['incorrect_red']

    left_sub = slide2.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.3))
    tf = left_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "No guardrails --> scattered effort, wasted resources"
    p.font.size = Pt(10)
    p.font.italic = True

    # Empty strategy box
    empty_strat = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.2), Inches(0.5))
    empty_strat.fill.solid()
    empty_strat.fill.fore_color.rgb = COLORS['white']
    empty_strat.line.dash_style = 2  # Dashed
    tf = empty_strat.text_frame
    p = tf.paragraphs[0]
    p.text = "? No Strategy Defined"
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER

    # Goal box
    goal_left = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(5.2), Inches(0.4))
    goal_left.fill.solid()
    goal_left.fill.fore_color.rgb = COLORS['strategy_blue']
    tf = goal_left.text_frame
    p = tf.paragraphs[0]
    p.text = "GOAL: Grow Revenue"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # No guardrails label
    no_guard = slide2.shapes.add_textbox(Inches(2.0), Inches(3.0), Inches(2.5), Inches(0.3))
    tf = no_guard.text_frame
    p = tf.paragraphs[0]
    p.text = "NO GUARDRAILS"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['incorrect_red']
    p.alignment = PP_ALIGN.CENTER

    # Scattered people/tactics description
    scatter_text = slide2.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(5.2), Inches(2.0))
    tf = scatter_text.text_frame
    tf.word_wrap = True
    items = [
        "Employee A --> \"Cut all costs\"",
        "Employee B --> \"Chase enterprise accounts\"",
        "Employee C --> \"Expand globally\"",
        "Employee D --> \"Discount heavily\"",
        "Employee E --> \"Build new features\"",
        "Employee F --> \"Focus on SMB\""
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(9)

    # Result box (negative)
    result_left = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.6), Inches(5.2), Inches(0.45))
    result_left.fill.solid()
    result_left.fill.fore_color.rgb = COLORS['incorrect_red']
    tf = result_left.text_frame
    p = tf.paragraphs[0]
    p.text = "RESULT: Conflicting priorities, organizational drift"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Right panel: WITH STRATEGY
    right_panel = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.1), Inches(6.2), Inches(5.2))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(0xd1, 0xfa, 0xe5)  # Light green
    right_panel.line.color.rgb = COLORS['goals_green']

    # Right header
    right_header = slide2.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = "WITH STRATEGY"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['goals_green']

    right_sub = slide2.shapes.add_textbox(Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.3))
    tf = right_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Clear guardrails --> aligned execution, compounding advantage"
    p.font.size = Pt(10)
    p.font.italic = True

    # Strategy box with trade-off
    strat_right = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(1.9), Inches(5.2), Inches(0.7))
    strat_right.fill.solid()
    strat_right.fill.fore_color.rgb = COLORS['goals_green']
    tf = strat_right.text_frame
    p = tf.paragraphs[0]
    p.text = "STRATEGY: Win mid-market manufacturing"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Trade-off: NOT enterprise, NOT horizontal"
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER

    # Goal box
    goal_right = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(2.7), Inches(5.2), Inches(0.4))
    goal_right.fill.solid()
    goal_right.fill.fore_color.rgb = COLORS['strategy_blue']
    tf = goal_right.text_frame
    p = tf.paragraphs[0]
    p.text = "GOAL: 500 Mfg Customers by 2026"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Guardrails label left
    guard_left = slide2.shapes.add_textbox(Inches(7.1), Inches(3.3), Inches(0.4), Inches(2.0))
    tf = guard_left.text_frame
    p = tf.paragraphs[0]
    p.text = "G\nU\nA\nR\nD\nR\nA\nI\nL\nS"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS['goals_green']

    # Guardrails label right
    guard_right = slide2.shapes.add_textbox(Inches(12.3), Inches(3.3), Inches(0.4), Inches(2.0))
    tf = guard_right.text_frame
    p = tf.paragraphs[0]
    p.text = "G\nU\nA\nR\nD\nR\nA\nI\nL\nS"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = COLORS['goals_green']

    # Aligned people/tactics description
    align_text = slide2.shapes.add_textbox(Inches(7.5), Inches(3.2), Inches(4.8), Inches(2.0))
    tf = align_text.text_frame
    tf.word_wrap = True
    items = [
        "All employees aligned:",
        "  --> NAM conference sponsor",
        "  --> 3 vertical sales reps",
        "  --> Mfg LinkedIn campaign",
        "  --> Industry case studies",
        "  --> Plant manager outreach"
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(9)

    # Result box (positive)
    result_right = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(5.6), Inches(5.2), Inches(0.45))
    result_right.fill.solid()
    result_right.fill.fore_color.rgb = COLORS['goals_green']
    tf = result_right.text_frame
    p = tf.paragraphs[0]
    p.text = "RESULT: Reinforcing tactics, compounding advantage"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Key insight footer
    insight_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.4), Inches(13.333), Inches(0.5))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = COLORS['dark_gray']
    insight_box.line.fill.background()
    tf = insight_box.text_frame
    p = tf.paragraphs[0]
    p.text = '"The essence of strategy is choosing what not to do." -- Michael Porter'
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Trade-offs create guardrails that align everyone's tactics toward the goal."
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLORS['white']
    p2.alignment = PP_ALIGN.CENTER

    # Sources
    source_box = slide2.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12), Inches(0.3))
    tf = source_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'Porter, M. "What Is Strategy?" HBR, 1996 | McKinsey Quarterly. "Strategic choices that lead to results"'
    p.font.size = Pt(8)
    p.font.color.rgb = COLORS['dark_gray']

    # Save
    pptx_path = os.path.join(OUTPUT_DIR, "Strategy_Slides.pptx")
    prs.save(pptx_path)
    print(f"PowerPoint saved: {pptx_path}")
    return pptx_path


def create_word_slide1():
    """Create detailed Word document for Slide 1."""
    doc = Document()

    # Title
    title = doc.add_heading('Strategy vs. Goals vs. Tactics', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph('Understanding the Hierarchy of Corporate Planning')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    # Introduction
    doc.add_heading('Introduction', level=1)
    doc.add_paragraph(
        'One of the most common sources of organizational confusion is the conflation of strategy, goals, and tactics. '
        'While these terms are often used interchangeably in business conversations, they represent fundamentally different '
        'concepts in the hierarchy of corporate planning. Understanding these distinctions is essential for effective '
        'strategic management and organizational alignment.'
    )
    doc.add_paragraph(
        'This document provides a comprehensive breakdown of each concept, including definitions, key components, '
        'and illustrative examples to help distinguish between them.'
    )

    # Strategy Section
    doc.add_heading('STRATEGY: Where & How to Compete', level=1)

    doc.add_heading('Definition', level=2)
    doc.add_paragraph(
        'An integrated set of choices that uniquely position an organization to create sustainable competitive advantage--'
        'including explicit trade-offs about what NOT to do.'
    )

    doc.add_heading('Key Components', level=2)
    components = [
        ('Where to Play', 'Markets, customers, geographies, and product categories the organization will pursue. '
         'This defines the competitive arena and target audience.'),
        ('How to Win', 'The unique value proposition and competitive advantages that will enable success in chosen markets. '
         'This explains why customers will choose you over alternatives.'),
        ('Trade-offs', 'Explicit choices about what NOT to do. Strategy requires sacrifice--attempting to be everything '
         'to everyone results in competitive mediocrity.'),
        ('Activity Fit', 'Ensuring all organizational activities reinforce each other and the strategic position. '
         'Individual activities should compound rather than conflict.')
    ]
    for comp, desc in components:
        p = doc.add_paragraph()
        p.add_run(f'{comp}: ').bold = True
        p.add_run(desc)

    doc.add_heading('Correct Examples', level=2)
    correct = [
        ('"Dominate mid-market manufacturing software through vertical specialization. We will NOT pursue enterprise accounts or horizontal applications."',
         'Clear where to play (mid-market manufacturing), how to win (vertical specialization), and explicit trade-offs (no enterprise, no horizontal).'),
        ('"Win in short-haul, point-to-point routes with no meals, no seat assignments, no hub transfers--enabling rapid turnaround and lowest costs." (Southwest Airlines)',
         'Specific market positioning with explicit operational trade-offs that reinforce each other.'),
        ('"Serve young, price-conscious furniture buyers through self-service and flat-pack design--sacrificing customers wanting delivery and assembly." (IKEA)',
         'Clear target customer, unique value proposition, and explicit sacrifice of a customer segment.')
    ]
    for example, explanation in correct:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this works: {explanation}')

    doc.add_heading('Incorrect Examples', level=2)
    incorrect = [
        ('"Be the best software company and grow revenue 20%."',
         'No choices, no trade-offs--this is a goal disguised as strategy. It lacks positioning and differentiation.'),
        ('"Delight our customers and be number one in the market."',
         'Vague aspiration with no positioning, no differentiation, and no guidance on resource allocation.'),
        ('"Focus on innovation and quality."',
         'No specific choices about where to compete or what to sacrifice. Every competitor could claim the same.')
    ]
    for example, explanation in incorrect:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this fails: {explanation}')

    # Goals Section
    doc.add_heading('GOALS: What Outcomes to Achieve', level=1)

    doc.add_heading('Definition', level=2)
    doc.add_paragraph(
        'Specific, measurable outcomes with defined timeframes that translate strategic intent into concrete targets.'
    )

    doc.add_heading('Key Components (SMART Framework)', level=2)
    components = [
        ('Specific', 'Clear articulation of what success looks like. Ambiguity is the enemy of execution.'),
        ('Measurable', 'Quantifiable metrics to track progress. If you cannot measure it, you cannot manage it.'),
        ('Achievable', 'Realistic given available resources, capabilities, and market conditions.'),
        ('Relevant', 'Aligned with strategic objectives. Goals should advance the strategy, not distract from it.'),
        ('Time-bound', 'Defined deadline for achievement. Open-ended goals lose urgency and accountability.')
    ]
    for comp, desc in components:
        p = doc.add_paragraph()
        p.add_run(f'{comp}: ').bold = True
        p.add_run(desc)

    doc.add_heading('Correct Examples', level=2)
    correct = [
        ('"Acquire 500 new manufacturing customers in North America by December 2026."',
         'Specific (manufacturing customers, North America), measurable (500), time-bound (December 2026).'),
        ('"Increase recurring subscription revenue to 30% of total revenue by FY2026."',
         'Clear metric, specific target, defined timeframe.'),
        ('"Achieve #2 market share position in European market within three years."',
         'Measurable position, specific geography, clear timeline.'),
        ('"Reduce customer acquisition cost by 25% while maintaining conversion rate by Q4 2025."',
         'Quantified improvement with guard rail (maintaining conversion) and deadline.')
    ]
    for example, explanation in correct:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this works: {explanation}')

    doc.add_heading('Incorrect Examples', level=2)
    incorrect = [
        ('"Grow market share significantly."',
         'Vague, no metrics, no deadline--unmeasurable and unaccountable.'),
        ('"Become more profitable."',
         'No specific target or timeframe. How much more? By when?'),
        ('"Improve customer satisfaction."',
         'No quantifiable metric, no deadline, no way to know when success is achieved.')
    ]
    for example, explanation in incorrect:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this fails: {explanation}')

    # Tactics Section
    doc.add_heading('TACTICS: How to Execute', level=1)

    doc.add_heading('Definition', level=2)
    doc.add_paragraph(
        'Specific actions, initiatives, or methods employed to execute strategy and achieve goals--operational, '
        'short-term, and adjustable.'
    )

    doc.add_heading('Key Components', level=2)
    components = [
        ('Short-term', 'Days to months timeframe. Tactics are the near-term activities that drive progress.'),
        ('Actionable', 'Concrete activities with clear steps. Someone should be able to start immediately.'),
        ('Adjustable', 'Can be changed without altering strategy. If a tactic fails, try another one.'),
        ('Owned', 'Clear accountability and responsibility. Every tactic needs an owner.'),
        ('Measurable', 'Trackable outputs and milestones. Progress should be visible.')
    ]
    for comp, desc in components:
        p = doc.add_paragraph()
        p.add_run(f'{comp}: ').bold = True
        p.add_run(desc)

    doc.add_heading('Correct Examples', level=2)
    correct = [
        ('"Sponsor NAM conference; hire 3 vertical sales reps; launch LinkedIn campaign targeting plant managers."',
         'Specific activities with clear actions that can begin immediately.'),
        ('"Launch referral program offering 20% discount for new customer acquisitions."',
         'Concrete program with defined incentive structure.'),
        ('"Redesign checkout flow to reduce cart abandonment by Q2."',
         'Specific improvement initiative with measurable outcome and deadline.'),
        ('"Run targeted LinkedIn campaign to plant operations managers in Midwest."',
         'Defined channel, audience, and geography.')
    ]
    for example, explanation in correct:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this works: {explanation}')

    doc.add_heading('Incorrect Examples', level=2)
    incorrect = [
        ('"Do more marketing and sales."',
         'Not actionable--no specific activity or owner. What marketing? What sales activities?'),
        ('"Improve our processes."',
         'Vague--no concrete action or timeline. Which processes? How?'),
        ('"Be more customer-focused."',
         'Aspirational statement, not an executable action. What specific behaviors or activities?')
    ]
    for example, explanation in incorrect:
        p = doc.add_paragraph(example, style='Quote')
        doc.add_paragraph(f'Why this fails: {explanation}')

    # Key Insight
    doc.add_heading('Key Insight', level=1)
    p = doc.add_paragraph()
    p.add_run('Strategy defines what NOT to do').bold = True
    p.add_run(' | ')
    p.add_run('Goals are measurable targets').bold = True
    p.add_run(' | ')
    p.add_run('Tactics are swappable actions').bold = True

    doc.add_paragraph(
        'The most important distinction is that strategy involves trade-offs--explicit choices about what you will NOT do. '
        'Without trade-offs, you do not have a strategy, you have a goal. Tactics, on the other hand, are interchangeable; '
        'if one does not work, you can try another without changing your strategy or goals.'
    )

    # Sources
    doc.add_heading('Sources', level=1)
    sources = [
        'Porter, M. "What Is Strategy?" Harvard Business Review, 1996',
        'Martin, R. "The Big Lie of Strategic Planning" Harvard Business Review, 2014',
        'McKinsey & Company. "Strategy to beat the odds"'
    ]
    for source in sources:
        doc.add_paragraph(source, style='List Bullet')

    # Save
    docx_path = os.path.join(OUTPUT_DIR, "Slide1_Strategy_vs_Goals_vs_Tactics_Detailed.docx")
    doc.save(docx_path)
    print(f"Word document saved: {docx_path}")
    return docx_path


def create_word_slide2():
    """Create detailed Word document for Slide 2."""
    doc = Document()

    # Title
    title = doc.add_heading('Strategy Creates Alignment--Its Absence Creates Chaos', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    # Introduction
    doc.add_heading('Introduction', level=1)
    doc.add_paragraph(
        'The primary function of strategy is not to predict the future or to optimize operations--it is to create alignment. '
        'When an organization has a clear strategy with explicit trade-offs, every employee can make decisions that move '
        'the organization in the same direction. Without strategy, well-intentioned employees pursue conflicting initiatives, '
        'wasting resources and creating organizational drift.'
    )
    doc.add_paragraph(
        'This document illustrates the contrast between organizations with and without clear strategic direction, '
        'demonstrating how trade-offs serve as guardrails that channel effort toward a common goal.'
    )

    # Without Strategy Section
    doc.add_heading('WITHOUT STRATEGY: The Problem of Misalignment', level=1)

    doc.add_heading('The Scenario', level=2)
    doc.add_paragraph(
        'Consider an organization with a goal to "grow revenue" but no clear strategy defining how to compete or what trade-offs to make.'
    )

    doc.add_heading('What Happens', level=2)
    doc.add_paragraph('Without strategic guardrails, different employees pursue different paths to the same goal:')

    misaligned = [
        ('Finance Team', 'Focuses on cutting costs to improve margins, potentially sacrificing growth investments'),
        ('Sales Team A', 'Chases large enterprise accounts that require long sales cycles and heavy customization'),
        ('Sales Team B', 'Pursues global expansion into new markets with different requirements'),
        ('Marketing Team', 'Implements heavy discounting to drive volume, eroding brand value'),
        ('Product Team', 'Builds features for whatever customer is loudest, creating a bloated product'),
        ('Operations Team', 'Optimizes for cost reduction, potentially sacrificing quality or speed')
    ]
    for team, action in misaligned:
        p = doc.add_paragraph()
        p.add_run(f'{team}: ').bold = True
        p.add_run(action)

    doc.add_heading('The Result', level=2)
    results = [
        'Conflicting priorities across departments',
        'Resource allocation battles with no clear resolution criteria',
        'Organizational drift as different parts of the company pull in different directions',
        'Employee frustration and disengagement',
        'Wasted resources on initiatives that cancel each other out',
        'Loss of competitive position as focus is diluted'
    ]
    for result in results:
        doc.add_paragraph(result, style='List Bullet')

    doc.add_heading('Why This Happens', level=2)
    doc.add_paragraph(
        'When an organization lacks a clear strategy, employees must guess at what is most important. Each employee, '
        'operating with their own perspective and incentives, makes reasonable choices that may conflict with others. '
        'The goal "grow revenue" provides no guidance on HOW to grow revenue or what NOT to do in pursuit of that goal.'
    )

    # With Strategy Section
    doc.add_heading('WITH STRATEGY: The Power of Guardrails', level=1)

    doc.add_heading('The Scenario', level=2)
    doc.add_paragraph(
        'Now consider the same organization with a clear strategy: "Win in mid-market manufacturing software through '
        'vertical specialization. We will NOT pursue enterprise accounts. We will NOT build horizontal applications."'
    )
    doc.add_paragraph(
        'The goal remains measurable: "Acquire 500 new manufacturing customers in North America by December 2026."'
    )

    doc.add_heading('What Changes', level=2)
    doc.add_paragraph('The strategy creates guardrails that align all employees toward the same objective:')

    aligned = [
        ('Sales Team', 'Focuses exclusively on mid-market manufacturing companies, developing deep vertical expertise'),
        ('Marketing Team', 'Sponsors NAM (National Association of Manufacturers) conference, builds manufacturing case studies'),
        ('Product Team', 'Builds features specifically for manufacturing workflows, ignoring requests from other verticals'),
        ('Finance Team', 'Allocates resources to manufacturing-focused initiatives, deprioritizes others'),
        ('HR Team', 'Hires sales reps with manufacturing industry background'),
        ('Customer Success', 'Develops manufacturing-specific onboarding and support playbooks')
    ]
    for team, action in aligned:
        p = doc.add_paragraph()
        p.add_run(f'{team}: ').bold = True
        p.add_run(action)

    doc.add_heading('Example Aligned Tactics', level=2)
    tactics = [
        'Sponsor NAM annual conference (brand visibility to target market)',
        'Hire 3 vertical sales representatives with manufacturing industry experience',
        'Launch LinkedIn campaign targeting plant operations managers',
        'Develop manufacturing-specific case studies and ROI calculators',
        'Build integrations with popular manufacturing ERP systems',
        'Create plant manager webinar series on digital transformation'
    ]
    for tactic in tactics:
        doc.add_paragraph(tactic, style='List Bullet')

    doc.add_heading('The Result', level=2)
    results = [
        'Reinforcing tactics that compound effectiveness',
        'Clear decision-making criteria ("Does this serve mid-market manufacturing?")',
        'Efficient resource allocation with minimal conflict',
        'Deep expertise development in the target market',
        'Strong brand recognition within the vertical',
        'Compounding competitive advantage over time'
    ]
    for result in results:
        doc.add_paragraph(result, style='List Bullet')

    # The Role of Trade-offs
    doc.add_heading('The Role of Trade-offs as Guardrails', level=1)

    doc.add_paragraph(
        'The magic of strategy lies not in what you choose to do, but in what you choose NOT to do. '
        'Trade-offs create guardrails that channel organizational energy.'
    )

    doc.add_heading('How Trade-offs Function', level=2)
    tradeoffs = [
        ('They eliminate options', 'By saying "NOT enterprise," the strategy eliminates a class of decisions. '
         'Sales does not need to debate whether to pursue a Fortune 500 opportunity--the answer is already decided.'),
        ('They resolve conflicts', 'When product and sales disagree about prioritization, the strategy provides the answer. '
         'If it does not serve mid-market manufacturing, it is deprioritized.'),
        ('They enable speed', 'Decisions that would require executive deliberation can be made by front-line employees. '
         'They know the guardrails.'),
        ('They build expertise', 'By constraining focus, the organization develops deep capability in its chosen area, '
         'creating sustainable competitive advantage.')
    ]
    for point, explanation in tradeoffs:
        p = doc.add_paragraph()
        p.add_run(f'{point}: ').bold = True
        p.add_run(explanation)

    # Key Insight
    doc.add_heading('Key Insight', level=1)
    p = doc.add_paragraph()
    p.add_run('"The essence of strategy is choosing what not to do."').bold = True
    p.add_run(' -- Michael Porter')

    doc.add_paragraph(
        'Trade-offs create guardrails that align everyone\'s tactics toward the goal. Without trade-offs, you have a goal, not a strategy. '
        'Without guardrails, you have motion without direction.'
    )

    # Practical Application
    doc.add_heading('Practical Application', level=1)

    doc.add_heading('Questions to Test Your Strategy', level=2)
    questions = [
        'Does your strategy explicitly state what you will NOT do?',
        'Can a front-line employee use the strategy to make decisions without executive input?',
        'Would your strategy force you to turn away certain customers or opportunities?',
        'Are all major initiatives clearly connected to the strategic positioning?',
        'Do your trade-offs create sustainable competitive advantage?'
    ]
    for q in questions:
        doc.add_paragraph(q, style='List Bullet')

    doc.add_heading('Warning Signs of Missing Strategy', level=2)
    warnings = [
        'Every opportunity is evaluated "on its merits" without clear criteria',
        'Different departments pursue conflicting priorities',
        'Resource allocation is decided by politics rather than strategy',
        'The organization is "good at many things" but dominant at none',
        'Employees cannot articulate what makes the organization different'
    ]
    for w in warnings:
        doc.add_paragraph(w, style='List Bullet')

    # Sources
    doc.add_heading('Sources', level=1)
    sources = [
        'Porter, M. "What Is Strategy?" Harvard Business Review, 1996',
        'McKinsey Quarterly. "Strategic choices that lead to results"',
        'Lafley, A.G. and Martin, R. "Playing to Win: How Strategy Really Works" 2013'
    ]
    for source in sources:
        doc.add_paragraph(source, style='List Bullet')

    # Save
    docx_path = os.path.join(OUTPUT_DIR, "Slide2_Strategy_Creates_Alignment_Detailed.docx")
    doc.save(docx_path)
    print(f"Word document saved: {docx_path}")
    return docx_path


if __name__ == "__main__":
    print("Exporting Strategy Slides...")
    print("=" * 50)

    pptx_path = create_powerpoint()
    print()

    docx1_path = create_word_slide1()
    print()

    docx2_path = create_word_slide2()
    print()

    print("=" * 50)
    print("Export complete!")
    print(f"\nFiles created in: {OUTPUT_DIR}")
    print(f"  1. {os.path.basename(pptx_path)}")
    print(f"  2. {os.path.basename(docx1_path)}")
    print(f"  3. {os.path.basename(docx2_path)}")
