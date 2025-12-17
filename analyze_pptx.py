#!/usr/bin/env python3
"""
Comprehensive PowerPoint Analysis Tool
Extracts structure, content, design elements, and visual frameworks
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from pathlib import Path

def analyze_color(color):
    """Extract color information"""
    try:
        if color.type == 1:  # RGB
            return {
                'type': 'RGB',
                'rgb': f'#{color.rgb:06X}' if hasattr(color, 'rgb') else None
            }
        elif color.type == 2:  # Theme
            return {'type': 'Theme', 'theme_color': str(color.theme_color)}
        else:
            return {'type': 'Other'}
    except:
        return None

def analyze_font(font):
    """Extract font information"""
    info = {}
    try:
        if font.name:
            info['name'] = font.name
        if font.size:
            info['size'] = font.size.pt
        if font.bold:
            info['bold'] = True
        if font.italic:
            info['italic'] = True
        if font.color:
            info['color'] = analyze_color(font.color)
    except:
        pass
    return info

def analyze_shape(shape, shape_idx):
    """Analyze individual shape properties"""
    shape_info = {
        'index': shape_idx,
        'type': str(shape.shape_type),
        'name': shape.name if hasattr(shape, 'name') else None,
    }

    # Position and size
    try:
        shape_info['position'] = {
            'left': shape.left.inches if hasattr(shape, 'left') else None,
            'top': shape.top.inches if hasattr(shape, 'top') else None,
            'width': shape.width.inches if hasattr(shape, 'width') else None,
            'height': shape.height.inches if hasattr(shape, 'height') else None,
        }
    except:
        pass

    # Text content
    if hasattr(shape, 'text') and shape.text:
        shape_info['text'] = shape.text

    # Text frame details
    if hasattr(shape, 'text_frame'):
        try:
            text_frame = shape.text_frame
            paragraphs = []
            for para in text_frame.paragraphs:
                para_info = {
                    'text': para.text,
                    'level': para.level,
                    'alignment': str(para.alignment) if para.alignment else None,
                }

                # Run details (formatting)
                runs = []
                for run in para.runs:
                    run_info = {
                        'text': run.text,
                        'font': analyze_font(run.font)
                    }
                    runs.append(run_info)

                if runs:
                    para_info['runs'] = runs

                paragraphs.append(para_info)

            if paragraphs:
                shape_info['paragraphs'] = paragraphs
        except Exception as e:
            shape_info['text_frame_error'] = str(e)

    # Fill color
    if hasattr(shape, 'fill'):
        try:
            if shape.fill.type == 1:  # Solid fill
                shape_info['fill'] = analyze_color(shape.fill.fore_color)
        except:
            pass

    # Line/border
    if hasattr(shape, 'line'):
        try:
            line_info = {}
            if shape.line.color:
                line_info['color'] = analyze_color(shape.line.color)
            if shape.line.width:
                line_info['width'] = shape.line.width.pt
            if line_info:
                shape_info['line'] = line_info
        except:
            pass

    # Table detection
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            table = shape.table
            table_data = []
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data.append(row_data)
            shape_info['table_data'] = table_data
        except Exception as e:
            shape_info['table_error'] = str(e)

    # Chart detection
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        shape_info['has_chart'] = True
        try:
            chart = shape.chart
            shape_info['chart_type'] = str(chart.chart_type)
        except:
            pass

    # Group detection
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shape_info['is_group'] = True
        try:
            grouped_shapes = []
            for sub_shape in shape.shapes:
                grouped_shapes.append(analyze_shape(sub_shape, len(grouped_shapes)))
            shape_info['grouped_shapes'] = grouped_shapes
        except:
            pass

    return shape_info

def analyze_slide(slide, slide_num):
    """Analyze a single slide comprehensively"""
    slide_info = {
        'number': slide_num,
        'slide_id': slide.slide_id if hasattr(slide, 'slide_id') else None,
    }

    # Slide layout
    try:
        slide_info['layout_name'] = slide.slide_layout.name
    except:
        pass

    # Title extraction
    if hasattr(slide, 'shapes') and slide.shapes.title:
        slide_info['title'] = slide.shapes.title.text

    # Notes
    if hasattr(slide, 'notes_slide'):
        try:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                slide_info['notes'] = notes_text
        except:
            pass

    # All shapes analysis
    shapes = []
    for shape_idx, shape in enumerate(slide.shapes):
        shapes.append(analyze_shape(shape, shape_idx))

    slide_info['shapes'] = shapes
    slide_info['shape_count'] = len(shapes)

    return slide_info

def analyze_presentation(pptx_path):
    """Main analysis function"""
    prs = Presentation(pptx_path)

    analysis = {
        'file_path': str(pptx_path),
        'metadata': {},
        'slide_dimensions': {},
        'slides': []
    }

    # Metadata
    try:
        core_props = prs.core_properties
        analysis['metadata'] = {
            'title': core_props.title,
            'author': core_props.author,
            'subject': core_props.subject,
            'created': str(core_props.created),
            'modified': str(core_props.modified),
        }
    except:
        pass

    # Slide dimensions
    analysis['slide_dimensions'] = {
        'width_inches': prs.slide_width.inches,
        'height_inches': prs.slide_height.inches,
    }

    # Total slide count
    analysis['total_slides'] = len(prs.slides)

    # Analyze each slide
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = analyze_slide(slide, slide_num)
        analysis['slides'].append(slide_info)

    return analysis

def generate_markdown_report(analysis):
    """Generate human-readable markdown report"""
    report = []
    report.append("# PowerPoint Analysis Report\n")
    report.append(f"**File:** {analysis['file_path']}\n")
    report.append(f"**Total Slides:** {analysis['total_slides']}\n")

    # Metadata
    if analysis['metadata']:
        report.append("\n## Metadata")
        for key, value in analysis['metadata'].items():
            if value:
                report.append(f"- **{key.title()}:** {value}")
        report.append("")

    # Dimensions
    dims = analysis['slide_dimensions']
    report.append(f"\n## Slide Dimensions")
    report.append(f"- Width: {dims['width_inches']:.2f} inches")
    report.append(f"- Height: {dims['height_inches']:.2f} inches")
    report.append("")

    # Slide-by-slide analysis
    report.append("\n## Slide-by-Slide Analysis\n")

    for slide in analysis['slides']:
        report.append(f"### Slide {slide['number']}")

        if 'title' in slide:
            report.append(f"**Title:** {slide['title']}\n")

        if 'layout_name' in slide:
            report.append(f"**Layout:** {slide['layout_name']}")

        report.append(f"**Shape Count:** {slide['shape_count']}\n")

        # Content extraction
        report.append("#### Content:\n")
        for shape in slide['shapes']:
            if 'text' in shape and shape['text'].strip():
                # Determine if it's likely a title or body text
                shape_type = shape.get('name', 'Unknown')
                text = shape['text'].strip()

                if 'title' in shape_type.lower():
                    report.append(f"**[Title]** {text}\n")
                elif 'subtitle' in shape_type.lower():
                    report.append(f"**[Subtitle]** {text}\n")
                else:
                    # Check for bullet points
                    if 'paragraphs' in shape:
                        report.append("**[Text Box/Body]**")
                        for para in shape['paragraphs']:
                            if para['text'].strip():
                                indent = "  " * para.get('level', 0)
                                report.append(f"{indent}- {para['text']}")
                        report.append("")
                    else:
                        report.append(f"**[Text]** {text}\n")

            # Table content
            if 'table_data' in shape:
                report.append("**[Table]**")
                for row in shape['table_data']:
                    report.append("| " + " | ".join(row) + " |")
                report.append("")

            # Chart notation
            if 'has_chart' in shape:
                chart_type = shape.get('chart_type', 'Unknown')
                report.append(f"**[Chart: {chart_type}]**\n")

        # Notes
        if 'notes' in slide:
            report.append(f"\n**Notes:**\n{slide['notes']}\n")

        report.append("\n---\n")

    return "\n".join(report)

def generate_design_analysis(analysis):
    """Generate design-focused analysis"""
    report = []
    report.append("# Design & Formatting Analysis\n")

    # Collect all fonts used
    fonts = set()
    colors = set()

    for slide in analysis['slides']:
        for shape in slide['shapes']:
            if 'paragraphs' in shape:
                for para in shape['paragraphs']:
                    if 'runs' in para:
                        for run in para['runs']:
                            if 'font' in run:
                                font_info = run['font']
                                if 'name' in font_info:
                                    fonts.add(font_info['name'])
                                if 'color' in font_info and font_info['color']:
                                    if 'rgb' in font_info['color']:
                                        colors.add(font_info['color']['rgb'])

            if 'fill' in shape and shape['fill']:
                if 'rgb' in shape['fill']:
                    colors.add(shape['fill']['rgb'])

    report.append("## Typography")
    report.append(f"**Fonts Used:** {', '.join(sorted(fonts)) if fonts else 'None detected'}\n")

    report.append("## Color Palette")
    report.append(f"**Colors Used:** {', '.join(sorted(colors)) if colors else 'None detected'}\n")

    report.append("## Slide Structure Summary\n")
    for slide in analysis['slides']:
        title = slide.get('title', f"Slide {slide['number']} (No Title)")
        layout = slide.get('layout_name', 'Unknown')
        report.append(f"- **{title}** [{layout}] - {slide['shape_count']} shapes")

    return "\n".join(report)

def generate_recommendations(analysis):
    """Generate recommendations for executive-ready slides"""
    report = []
    report.append("# Executive Readiness Assessment & Recommendations\n")

    # Analyze current state
    avg_shapes = sum(s['shape_count'] for s in analysis['slides']) / len(analysis['slides'])
    slides_with_titles = sum(1 for s in analysis['slides'] if 'title' in s)

    report.append("## Current State Analysis\n")
    report.append(f"- **Total Slides:** {analysis['total_slides']}")
    report.append(f"- **Average Shapes per Slide:** {avg_shapes:.1f}")
    report.append(f"- **Slides with Titles:** {slides_with_titles}/{analysis['total_slides']}")
    report.append("")

    report.append("## Strengths to Preserve\n")

    # Identify strengths
    if slides_with_titles == analysis['total_slides']:
        report.append("- All slides have clear titles")

    # Check for consistent structure
    layouts = [s.get('layout_name') for s in analysis['slides']]
    unique_layouts = set(layouts)
    if len(unique_layouts) <= 3:
        report.append("- Consistent slide layout usage")

    report.append("")
    report.append("## Areas for Improvement\n")

    # Recommendations
    if avg_shapes > 10:
        report.append("- **Simplify:** Average of {:.1f} shapes per slide may be too cluttered for executives".format(avg_shapes))
        report.append("  - Recommendation: Aim for 5-7 key elements per slide")

    if slides_with_titles < analysis['total_slides']:
        report.append(f"- **Add Titles:** {analysis['total_slides'] - slides_with_titles} slides missing titles")

    # Text density check
    text_heavy_slides = []
    for slide in analysis['slides']:
        total_chars = 0
        for shape in slide['shapes']:
            if 'text' in shape:
                total_chars += len(shape['text'])
        if total_chars > 500:
            text_heavy_slides.append(slide['number'])

    if text_heavy_slides:
        report.append(f"- **Reduce Text:** Slides {', '.join(map(str, text_heavy_slides))} appear text-heavy")
        report.append("  - Recommendation: Use bullet points, limit to 6 lines per slide")

    report.append("")
    report.append("## Executive Design Recommendations\n")
    report.append("1. **Visual Hierarchy:** Use size and color to guide attention")
    report.append("2. **White Space:** Increase margins and spacing between elements")
    report.append("3. **Data Visualization:** Replace text with charts/diagrams where possible")
    report.append("4. **Consistency:** Use consistent fonts, colors, and layouts throughout")
    report.append("5. **One Message per Slide:** Each slide should convey one key insight")
    report.append("6. **Professional Color Scheme:** Use corporate colors or professional palette")
    report.append("7. **Clear Call-to-Actions:** End sections with clear next steps")

    return "\n".join(report)

if __name__ == "__main__":
    pptx_path = "/Users/nicholaspate/Documents/01_Active/Corp_Strat/ezAI/insights/strategy-goals-tactics-slides.pptx"

    print("Analyzing PowerPoint presentation...")
    analysis = analyze_presentation(pptx_path)

    # Save detailed JSON
    json_path = Path(pptx_path).parent / "pptx_analysis_detailed.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(analysis, f, indent=2, ensure_ascii=False)
    print(f"Saved detailed analysis to: {json_path}")

    # Generate markdown report
    markdown_report = generate_markdown_report(analysis)
    report_path = Path(pptx_path).parent / "pptx_analysis_report.md"
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write(markdown_report)
    print(f"Saved content report to: {report_path}")

    # Generate design analysis
    design_report = generate_design_analysis(analysis)
    design_path = Path(pptx_path).parent / "pptx_design_analysis.md"
    with open(design_path, 'w', encoding='utf-8') as f:
        f.write(design_report)
    print(f"Saved design analysis to: {design_path}")

    # Generate recommendations
    recommendations = generate_recommendations(analysis)
    rec_path = Path(pptx_path).parent / "pptx_recommendations.md"
    with open(rec_path, 'w', encoding='utf-8') as f:
        f.write(recommendations)
    print(f"Saved recommendations to: {rec_path}")

    print("\nAnalysis complete!")
    print(f"Total slides analyzed: {analysis['total_slides']}")
